import tkinter as tk
from tkinter import ttk, messagebox
import threading
import sqlite3
import os
import re
import urllib.parse
import requests
import time
import ollama
from pptx import Presentation
from pptx.util import Pt, Inches
from pptx.dml.color import RGBColor
from dotenv import load_dotenv

load_dotenv()

# ─── MODELS ──────────────────────────────────────────────────────
LOCAL_MODELS = [
    "gemma:2b",
    "gemma4:latest",
]

UNSPLASH_KEY = os.getenv("UNSPLASH_API_KEY", "")

# ─── COLOURS ─────────────────────────────────────────────────────
BG_COLOR     = RGBColor(0x06, 0x5A, 0x82)
ACCENT_COLOR = RGBColor(0x02, 0xC3, 0x9A)
TEXT_LIGHT   = RGBColor(0xFF, 0xFF, 0xFF)
TEXT_BODY    = RGBColor(0xE8, 0xF4, 0xF8)
SLIDE_W      = Inches(13.33)
SLIDE_H      = Inches(7.5)

STOPWORDS = {
    "a","an","the","of","in","on","at","to","for","with","and","or",
    "is","are","was","were","be","been","by","from","this","that","as",
    "its","it","their","about","using","through","between","across"
}

# ─── STEP 1: Get slide titles outline from model ──────────────────
def get_slide_outline(client, model, topic, num_slides):
    """Ask model for a numbered list of slide titles only — very easy for small models."""
    prompt = (
        f'Create an outline for a {num_slides}-slide presentation on: "{topic}".\n'
        f'Reply with ONLY a numbered list of {num_slides} slide titles, one per line.\n'
        f'Example format:\n'
        f'1. Introduction to the Topic\n'
        f'2. Key Concepts\n'
        f'3. How It Works\n'
        f'Do not add any other text.'
    )
    resp = client.generate(model=model, prompt=prompt)
    raw  = getattr(resp, "response", str(resp)).strip()

    titles = []
    for line in raw.splitlines():
        line = line.strip()
        # Match "1. Title" or "1) Title" or "- Title"
        m = re.match(r'^[\d]+[.)]\s*(.+)', line)
        if m:
            title = re.sub(r'[*#_]', '', m.group(1)).strip()
            if title:
                titles.append(title)
        elif line.startswith('-') or line.startswith('•'):
            title = re.sub(r'^[-•]\s*', '', line).strip()
            title = re.sub(r'[*#_]', '', title).strip()
            if title:
                titles.append(title)

    # If parsing failed, generate generic titles
    if len(titles) < num_slides:
        base = ["Introduction", "Overview", "Key Concepts", "How It Works",
                "Applications", "Benefits", "Challenges", "Case Studies",
                "Future Trends", "Best Practices", "Tools & Technologies",
                "Getting Started", "Advanced Topics", "Real World Examples", "Conclusion"]
        while len(titles) < num_slides:
            titles.append(f"{base[len(titles) % len(base)]} - {topic}")

    return titles[:num_slides]


# ─── STEP 2: Get bullet content for one slide ────────────────────
def get_slide_content(client, model, topic, slide_title, slide_num, total):
    """Ask for bullets only — much simpler prompt that small models handle well."""
    prompt = (
        f'Presentation topic: "{topic}"\n'
        f'Slide {slide_num} of {total}: "{slide_title}"\n\n'
        f'Write exactly 4 bullet points for this slide.\n'
        f'Each bullet must:\n'
        f'- Start with a dash (-)\n'
        f'- Be one clear sentence\n'
        f'- Be specific to "{slide_title}"\n\n'
        f'Reply with ONLY the 4 bullet points, nothing else.'
    )
    resp = client.generate(model=model, prompt=prompt)
    raw  = getattr(resp, "response", str(resp)).strip()

    bullets = []
    for line in raw.splitlines():
        line = line.strip()
        if not line:
            continue
        # Strip any bullet marker
        clean = re.sub(r'^[-*•\d.)+]+\s*', '', line)
        clean = re.sub(r'\*\*(.*?)\*\*', r'\1', clean)
        clean = re.sub(r'\*(.*?)\*',     r'\1', clean)
        clean = re.sub(r'[#_]', '', clean).strip()
        if len(clean) > 10:   # ignore very short/empty lines
            bullets.append(clean)

    # Fallback if model returned nothing useful
    if not bullets:
        bullets = [
            f"Key aspect of {slide_title}",
            f"Important consideration in {slide_title}",
            f"How {slide_title} relates to {topic}",
            f"Practical applications of {slide_title}",
        ]

    return bullets[:6]


# ─── Image keywords: derived from topic + title, NOT the model ───
def make_image_keywords(topic, title):
    combined = f"{topic} {title}".lower()
    words    = re.findall(r'\b[a-z]{3,}\b', combined)
    unique   = list(dict.fromkeys(w for w in words if w not in STOPWORDS))
    return " ".join(unique[:4])


# ─── DB ──────────────────────────────────────────────────────────
def init_db(db_path):
    conn = sqlite3.connect(db_path)
    conn.execute('''CREATE TABLE IF NOT EXISTS slides (
        id INTEGER PRIMARY KEY AUTOINCREMENT,
        slide_number INTEGER, title TEXT, content TEXT,
        image_keywords TEXT, status TEXT DEFAULT "pending"
    )''')
    conn.commit(); conn.close()


# ─── PPTX BUILDER ────────────────────────────────────────────────
def clean_md(text):
    text = re.sub(r'^#{1,6}\s*', '', text, flags=re.MULTILINE)
    text = re.sub(r'\*\*(.*?)\*\*', r'\1', text)
    text = re.sub(r'\*(.*?)\*',     r'\1', text)
    text = re.sub(r'^[-*•]+\s*',   '',    text, flags=re.MULTILINE)
    return text.strip()

def add_bg(slide):
    bg = slide.shapes.add_shape(1, 0, 0, SLIDE_W, SLIDE_H)
    bg.fill.solid(); bg.fill.fore_color.rgb = BG_COLOR; bg.line.fill.background()
    sp = slide.shapes._spTree
    sp.remove(bg._element); sp.insert(2, bg._element)

def add_title_bar(slide, title):
    bar = slide.shapes.add_shape(1, 0, 0, SLIDE_W, Inches(1.1))
    bar.fill.solid(); bar.fill.fore_color.rgb = ACCENT_COLOR; bar.line.fill.background()
    txb = slide.shapes.add_textbox(Inches(0.4), Inches(0.1), Inches(12.5), Inches(0.9))
    tf  = txb.text_frame; tf.word_wrap = False
    p   = tf.paragraphs[0]; p.text = clean_md(title).upper()
    p.font.bold = True; p.font.size = Pt(26); p.font.name = "Calibri"
    p.font.color.rgb = TEXT_LIGHT

def add_content(slide, bullets):
    txb = slide.shapes.add_textbox(Inches(0.5), Inches(1.35), Inches(6.2), Inches(5.8))
    tf  = txb.text_frame; tf.word_wrap = True
    first = True
    for bullet in bullets[:6]:
        p = tf.paragraphs[0] if first else tf.add_paragraph()
        first = False
        run = p.add_run()
        run.text = "• " + clean_md(bullet)
        p.font.size = Pt(16); p.font.name = "Calibri"
        p.font.color.rgb = TEXT_BODY; p.space_after = Pt(8)

def add_image(slide, img_path):
    try:
        slide.shapes.add_picture(img_path, Inches(7.1), Inches(1.35), Inches(5.8), Inches(5.5))
    except Exception:
        pass

def build_pptx(db_path, output_path):
    conn  = sqlite3.connect(db_path)
    rows  = conn.execute("SELECT id, title, content FROM slides ORDER BY slide_number").fetchall()
    conn.close()
    prs   = Presentation()
    prs.slide_width = SLIDE_W; prs.slide_height = SLIDE_H
    blank = prs.slide_layouts[6]
    img_dir = os.path.join(os.path.dirname(db_path), "images")
    for db_id, title, content in rows:
        slide   = prs.slides.add_slide(blank)
        bullets = [l.strip() for l in (content or "").splitlines() if l.strip()]
        add_bg(slide)
        add_title_bar(slide, title or f"Slide {db_id}")
        add_content(slide, bullets)
        img = os.path.join(img_dir, f"slide_{db_id}.jpg")
        if os.path.exists(img):
            add_image(slide, img)
    prs.save(output_path)


# ─── IMAGE DOWNLOADER ────────────────────────────────────────────
def download_images(db_path, log):
    img_dir = os.path.join(os.path.dirname(db_path), "images")
    os.makedirs(img_dir, exist_ok=True)
    conn = sqlite3.connect(db_path)
    rows = conn.execute("SELECT id, image_keywords FROM slides ORDER BY slide_number").fetchall()
    conn.close()

    for db_id, keywords in rows:
        path = os.path.join(img_dir, f"slide_{db_id}.jpg")
        if os.path.exists(path) and os.path.getsize(path) > 1000:
            log(f"  🖼  Slide {db_id}: already exists, skipping")
            continue

        query = (keywords or "technology").strip()
        # Add page offset based on slide number so images aren't all the same
        page  = ((db_id - 1) % 5) + 1
        url   = (
            f"https://api.unsplash.com/search/photos"
            f"?query={urllib.parse.quote(query)}&per_page=5&page={page}"
            f"&orientation=landscape&client_id={UNSPLASH_KEY}"
        )
        try:
            r    = requests.get(url, timeout=10)
            data = r.json()
            if r.status_code == 200 and data.get("results"):
                # Pick result by slide_id index to vary images even for similar queries
                pick     = (db_id - 1) % len(data["results"])
                img_url  = data["results"][pick]["urls"]["regular"]
                img_data = requests.get(img_url, timeout=15).content
                with open(path, "wb") as f:
                    f.write(img_data)
                log(f"  🖼  Slide {db_id}: ✅ {query}")
            else:
                log(f"  ⚠  Slide {db_id}: no image ({r.status_code})")
            time.sleep(1.2)
        except Exception as e:
            log(f"  ⚠  Slide {db_id} error: {e}")


# ─── GUI ─────────────────────────────────────────────────────────
class App(tk.Tk):
    def __init__(self):
        super().__init__()
        self.title("AI PPT Generator")
        self.geometry("680x600")
        self.resizable(False, False)
        self.configure(bg="#0f172a")
        self._build_ui()

    def _build_ui(self):
        FONT      = ("Segoe UI", 11)
        FONT_HEAD = ("Segoe UI", 12, "bold")
        FG        = "#e2e8f0"
        ENTRY_BG  = "#1e293b"
        BTN_BG    = "#02c39a"
        BTN_FG    = "#0f172a"

        tk.Label(self, text="🎯  AI PPT Generator", font=("Segoe UI", 17, "bold"),
                 bg="#0f172a", fg=BTN_BG).pack(pady=(20, 4))
        tk.Label(self, text="Local Ollama Models", font=("Segoe UI", 9),
                 bg="#0f172a", fg="#64748b").pack(pady=(0, 16))

        frm = tk.Frame(self, bg="#0f172a")
        frm.pack(fill="x", padx=40)
        frm.columnconfigure(0, weight=3)
        frm.columnconfigure(1, weight=1)

        # Topic
        tk.Label(frm, text="Presentation Topic", font=FONT_HEAD, bg="#0f172a",
                 fg=FG, anchor="w").grid(row=0, column=0, columnspan=2, sticky="w", pady=(0,4))
        self.topic_var = tk.StringVar()
        tk.Entry(frm, textvariable=self.topic_var, font=FONT, bg=ENTRY_BG, fg=FG,
                 insertbackground=FG, relief="flat", bd=6).grid(
            row=1, column=0, columnspan=2, sticky="ew", pady=(0, 14))

        # Slides
        tk.Label(frm, text="Slides", font=FONT_HEAD, bg="#0f172a", fg=FG,
                 anchor="w").grid(row=2, column=0, sticky="w", pady=(0,4))
        tk.Label(frm, text="Model", font=FONT_HEAD, bg="#0f172a", fg=FG,
                 anchor="w").grid(row=2, column=1, sticky="w", padx=(14,0), pady=(0,4))

        self.slides_var = tk.IntVar(value=8)
        tk.Spinbox(frm, from_=3, to=20, textvariable=self.slides_var, font=FONT,
                   bg=ENTRY_BG, fg=FG, relief="flat", buttonbackground="#334155",
                   width=5).grid(row=3, column=0, sticky="w", pady=(0,14))

        self.model_var = tk.StringVar(value=LOCAL_MODELS[0])
        ttk.Combobox(frm, textvariable=self.model_var, values=LOCAL_MODELS,
                     font=FONT, state="readonly", width=22).grid(
            row=3, column=1, sticky="w", padx=(14,0), pady=(0,14))

        # Output folder
        tk.Label(frm, text="Output Folder", font=FONT_HEAD, bg="#0f172a",
                 fg=FG, anchor="w").grid(row=4, column=0, columnspan=2, sticky="w", pady=(0,4))
        self.out_var = tk.StringVar(value=os.path.expanduser("~/Desktop"))
        tk.Entry(frm, textvariable=self.out_var, font=FONT, bg=ENTRY_BG, fg=FG,
                 insertbackground=FG, relief="flat", bd=6).grid(
            row=5, column=0, columnspan=2, sticky="ew", pady=(0, 18))

        # Button
        self.btn = tk.Button(self, text="✨  Generate Presentation",
                             font=("Segoe UI", 12, "bold"),
                             bg=BTN_BG, fg=BTN_FG, relief="flat", bd=0,
                             padx=20, pady=10, cursor="hand2", command=self._start)
        self.btn.pack(pady=(0, 12))

        # Progress
        self.progress = ttk.Progressbar(self, length=600, mode="determinate")
        self.progress.pack(pady=(0, 8))

        # Log
        log_frame = tk.Frame(self, bg="#0f172a")
        log_frame.pack(fill="both", expand=True, padx=40, pady=(0, 20))
        self.log_box = tk.Text(log_frame, height=10, font=("Consolas", 9),
                               bg="#1e293b", fg="#94a3b8", relief="flat",
                               state="disabled", wrap="word")
        scroll = tk.Scrollbar(log_frame, command=self.log_box.yview)
        self.log_box.configure(yscrollcommand=scroll.set)
        self.log_box.pack(side="left", fill="both", expand=True)
        scroll.pack(side="right", fill="y")

        style = ttk.Style(self)
        style.theme_use("clam")
        style.configure("TCombobox", fieldbackground="#1e293b", background="#1e293b",
                        foreground="#e2e8f0", selectbackground="#1e293b")

    def _log(self, msg):
        self.log_box.configure(state="normal")
        self.log_box.insert("end", msg + "\n")
        self.log_box.see("end")
        self.log_box.configure(state="disabled")
        self.update_idletasks()

    def _start(self):
        topic = self.topic_var.get().strip()
        if not topic:
            messagebox.showwarning("Missing Topic", "Please enter a topic.")
            return
        self.btn.configure(state="disabled", text="Generating…")
        self.log_box.configure(state="normal"); self.log_box.delete("1.0", "end")
        self.log_box.configure(state="disabled")
        self.progress["value"] = 0
        threading.Thread(target=self._run, daemon=True).start()

    def _run(self):
        topic      = self.topic_var.get().strip()
        num_slides = self.slides_var.get()
        model      = self.model_var.get()
        out_folder = self.out_var.get().strip()
        os.makedirs(out_folder, exist_ok=True)

        safe   = re.sub(r'[^\w\s-]', '', topic).strip().replace(' ', '_')[:40]
        db_path    = os.path.join(out_folder, f"{safe}.db")
        pptx_path  = os.path.join(out_folder, f"{safe}.pptx")

        try:
            init_db(db_path)
            conn = sqlite3.connect(db_path)
            conn.execute("DELETE FROM slides"); conn.commit(); conn.close()

            self._log(f"📋 Topic  : {topic}")
            self._log(f"🤖 Model  : {model}")
            self._log(f"📊 Slides : {num_slides}\n")

            client      = ollama.Client(host="http://localhost:11434")
            total_steps = 1 + num_slides + num_slides + 1  # outline + content + images + build

            # ── Step 1: Outline ──────────────────────────────────
            self._log("📝 Getting slide outline…")
            titles = get_slide_outline(client, model, topic, num_slides)
            self._log(f"   Got {len(titles)} titles ✅\n")
            self.progress["value"] = (1 / total_steps) * 100

            # ── Step 2: Content per slide ────────────────────────
            for i, title in enumerate(titles, 1):
                self._log(f"✍  Slide {i}/{num_slides}: {title}")
                bullets  = get_slide_content(client, model, topic, title, i, num_slides)
                keywords = make_image_keywords(topic, title)

                conn = sqlite3.connect(db_path)
                conn.execute(
                    "INSERT INTO slides (slide_number,title,content,image_keywords,status) "
                    "VALUES (?,?,?,?,?)",
                    (i, title, "\n".join(bullets), keywords, "pending")
                )
                conn.commit(); conn.close()
                self._log(f"   ✅ {len(bullets)} bullets | 🔑 {keywords}")
                self.progress["value"] = ((1 + i) / total_steps) * 100

            # ── Step 3: Images ───────────────────────────────────
            self._log("\n🖼  Downloading images…")
            download_images(db_path, self._log)
            self.progress["value"] = ((1 + num_slides + num_slides) / total_steps) * 100

            # ── Step 4: Build PPTX ───────────────────────────────
            self._log("\n🔨 Building PPTX…")
            build_pptx(db_path, pptx_path)
            self.progress["value"] = 100

            self._log(f"\n✅ Done!\n   {pptx_path}")
            messagebox.showinfo("Done!", f"Saved to:\n{pptx_path}")
            os.startfile(pptx_path)

        except Exception as e:
            self._log(f"\n❌ Error: {e}")
            messagebox.showerror("Error", str(e))
        finally:
            self.btn.configure(state="normal", text="✨  Generate Presentation")


if __name__ == "__main__":
    App().mainloop()