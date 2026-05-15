import sqlite3
import os
import re
from pptx import Presentation
from pptx.util import Pt, Inches
from pptx.dml.color import RGBColor

# ── Colour palette ────────────────────────────────────────────────
BG_COLOR     = RGBColor(0x06, 0x5A, 0x82)   # deep blue
ACCENT_COLOR = RGBColor(0x02, 0xC3, 0x9A)   # mint
TEXT_LIGHT   = RGBColor(0xFF, 0xFF, 0xFF)
TEXT_BODY    = RGBColor(0xE8, 0xF4, 0xF8)   # ice white

SLIDE_W = Inches(13.33)
SLIDE_H = Inches(7.5)


def clean_markdown(text):
    """Strip markdown so raw AI responses still look clean on slides."""
    text = re.sub(r"^#{1,6}\s*", "", text, flags=re.MULTILINE)   # headings
    text = re.sub(r"\*\*(.*?)\*\*", r"\1", text)                  # bold
    text = re.sub(r"\*(.*?)\*",     r"\1", text)                  # italic
    text = re.sub(r"^[-*•]+\s*",   "",    text, flags=re.MULTILINE)  # bullets
    text = re.sub(r"^\*{3,}$",     "",    text, flags=re.MULTILINE)  # hr lines
    text = re.sub(r"\[([^\]]+)\]\([^)]+\)", r"\1", text)          # links
    return text.strip()


def extract_title_from_content(raw_content):
    """If populate_db stored the raw response, try to pull the TITLE: line."""
    m = re.search(r"TITLE:\s*(.+)", raw_content, re.IGNORECASE)
    if m:
        return clean_markdown(m.group(1)).strip()
    # Fall back to first heading
    m = re.search(r"^#{1,3}\s*(.+)", raw_content, re.MULTILINE)
    if m:
        return clean_markdown(m.group(1)).strip()
    return ""


def extract_bullets_from_content(raw_content):
    """Pull actual bullet lines, stripping everything else."""
    lines = []
    in_content = False
    for line in raw_content.splitlines():
        up = line.strip().upper()
        if up.startswith("CONTENT:") or up.startswith("BULLETS:"):
            in_content = True
            continue
        if up.startswith("IMAGE_PROMPT:") or up.startswith("TITLE:"):
            in_content = False
            continue
        if in_content:
            clean = re.sub(r"^[-*•]+\s*", "", line.strip())
            clean = re.sub(r"\*\*(.*?)\*\*", r"\1", clean)
            clean = re.sub(r"\*(.*?)\*",     r"\1", clean)
            if clean:
                lines.append(clean)

    # If old DB had no CONTENT: marker, fall back to cleaning every line
    if not lines:
        for line in raw_content.splitlines():
            clean = clean_markdown(line)
            if (clean
                    and not clean.upper().startswith("TITLE:")
                    and not clean.upper().startswith("IMAGE_PROMPT:")
                    and not clean.upper().startswith("CONTEXT NOTE")
                    and len(clean) > 5):
                lines.append(clean)

    return lines


def add_background(slide):
    bg = slide.shapes.add_shape(1, 0, 0, SLIDE_W, SLIDE_H)
    bg.fill.solid()
    bg.fill.fore_color.rgb = BG_COLOR
    bg.line.fill.background()
    sp_tree = slide.shapes._spTree
    sp_tree.remove(bg._element)
    sp_tree.insert(2, bg._element)


def add_title_bar(slide, title_text):
    bar = slide.shapes.add_shape(1, 0, 0, SLIDE_W, Inches(1.1))
    bar.fill.solid()
    bar.fill.fore_color.rgb = ACCENT_COLOR
    bar.line.fill.background()

    txb = slide.shapes.add_textbox(Inches(0.4), Inches(0.1), Inches(12.5), Inches(0.9))
    tf  = txb.text_frame
    tf.word_wrap = False
    p = tf.paragraphs[0]
    p.text           = title_text.upper()
    p.font.bold      = True
    p.font.size      = Pt(30)
    p.font.name      = "Calibri"
    p.font.color.rgb = TEXT_LIGHT


def add_content(slide, bullets):
    txb = slide.shapes.add_textbox(Inches(0.5), Inches(1.35), Inches(6.2), Inches(5.8))
    tf  = txb.text_frame
    tf.word_wrap = True

    first = True
    for bullet in bullets[:6]:   # cap at 6 bullets per slide
        p = tf.paragraphs[0] if first else tf.add_paragraph()
        first = False
        run = p.add_run()
        run.text         = "• " + bullet
        p.font.size      = Pt(16)
        p.font.name      = "Calibri"
        p.font.color.rgb = TEXT_BODY
        p.space_after    = Pt(8)


def add_image(slide, img_path):
    try:
        slide.shapes.add_picture(
            img_path,
            Inches(7.1), Inches(1.35),
            Inches(5.8), Inches(5.5)
        )
    except Exception as e:
        print(f"   ⚠ Image skipped: {e}")


def build_presentation(db_path="presentation.db", output="Final_Clean_Presentation.pptx"):
    if not os.path.exists(db_path):
        print(f"❌ {db_path} not found. Run populate_db.py first.")
        return

    conn   = sqlite3.connect(db_path)
    cursor = conn.cursor()

    try:
        cursor.execute("SELECT id, title, content FROM slides ORDER BY slide_number")
        rows = cursor.fetchall()
    except sqlite3.OperationalError as e:
        print(f"❌ DB error: {e}")
        return

    if not rows:
        print("⚠ No slides in database.")
        return

    prs              = Presentation()
    prs.slide_width  = SLIDE_W
    prs.slide_height = SLIDE_H
    blank_layout     = prs.slide_layouts[6]

    print(f"Building {len(rows)} slides...\n")

    for db_id, raw_title, raw_content in rows:
        slide = prs.slides.add_slide(blank_layout)

        # Determine best title
        title = (raw_title or "").strip()
        if not title or title.lower().startswith("slide "):
            title = extract_title_from_content(raw_content or "")
        if not title:
            title = f"Slide {db_id}"

        bullets = extract_bullets_from_content(raw_content or "")
        print(f"  Slide {db_id}: {title} ({len(bullets)} bullets)")

        add_background(slide)
        add_title_bar(slide, title)
        add_content(slide, bullets)

        img_path = f"images/slide_{db_id}.jpg"
        if os.path.exists(img_path):
            add_image(slide, img_path)

    try:
        prs.save(output)
        print(f"\n✅ Saved: {output}")
    except PermissionError:
        print(f"\n❌ Close {output} first, then re-run.")

    conn.close()


if __name__ == "__main__":
    build_presentation()
