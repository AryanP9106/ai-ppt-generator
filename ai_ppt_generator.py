"""
AI PPT Generator - Complete Solution
=====================================
Uses Ollama (local LLM) to generate slide content,
Unsplash API for relevant images, and python-pptx to assemble the deck.

Setup:
    pip install python-pptx requests ollama pillow

Usage:
    python ai_ppt_generator.py
    
    Then enter your topic and number of slides when prompted.
"""

import os
import re
import json
import time
import requests
import ollama
from pathlib import Path
from io import BytesIO
from PIL import Image
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

# ─────────────────────────────────────────────
# CONFIGURATION — Edit these
# ─────────────────────────────────────────────
UNSPLASH_ACCESS_KEY = "FBT7VO06efjln1SnF3l_yKjMYP5N7ReS1ti4qzjvig4"   # Get free key at unsplash.com/developers
OLLAMA_MODEL        = "gemma4:e4b"                      # Change to your local model (gemma3, llama3, mistral, etc.)
OUTPUT_FILE         = "AI_Presentation.pptx"
IMAGES_DIR          = Path("slide_images")
SLIDE_W             = Inches(13.33)   # Widescreen 16:9
SLIDE_H             = Inches(7.5)

# ─────────────────────────────────────────────
# COLOR PALETTE — Ocean Gradient theme
# ─────────────────────────────────────────────
CLR_BG_DARK   = RGBColor(0x06, 0x5A, 0x82)   # deep blue — title/section slides
CLR_BG_LIGHT  = RGBColor(0xFF, 0xFF, 0xFF)   # white — content slides
CLR_ACCENT    = RGBColor(0x02, 0xC3, 0x9A)   # mint green accent
CLR_TEXT_DARK = RGBColor(0x1A, 0x1A, 0x2E)   # near-black
CLR_TEXT_LITE = RGBColor(0xFF, 0xFF, 0xFF)   # white
CLR_MUTED     = RGBColor(0x55, 0x7A, 0x95)   # muted blue-grey


# ═══════════════════════════════════════════════════════════
# STEP 1 — LLM: Generate structured slide content
# ═══════════════════════════════════════════════════════════

def generate_slide_plan(topic: str, num_slides: int) -> list[dict]:
    """Ask LLM to produce a JSON array of slide objects."""
    print(f"\n🤖 Generating {num_slides}-slide plan for: '{topic}' ...")

    prompt = f"""You are a professional presentation designer.
Create a {num_slides}-slide presentation outline about: "{topic}"

Return ONLY a JSON array (no explanation, no markdown fences). Each element must have:
- "title": short slide title (max 8 words)
- "bullets": list of 3-4 concise bullet points (each max 12 words)
- "image_keyword": 1-3 word phrase for an Unsplash image search (specific, visual noun)
- "slide_type": one of "title", "content", "section", "conclusion"

Slide 1 must be type "title". Last slide must be type "conclusion".
Every 8-10 slides insert a "section" slide with a single-sentence summary as bullets[0].
All other slides are "content".

Example element:
{{"title": "Why Data Matters", "bullets": ["Data drives smarter decisions", "Real-time insights reduce costs", "Companies using data grow 23% faster"], "image_keyword": "data analytics dashboard", "slide_type": "content"}}

Now produce {num_slides} slides about "{topic}":"""

    response = ollama.chat(
        model=OLLAMA_MODEL,
        messages=[{"role": "user", "content": prompt}]
    )
    raw = response["message"]["content"].strip()

    # Strip markdown fences if model adds them
    raw = re.sub(r"^```(?:json)?\s*", "", raw)
    raw = re.sub(r"\s*```$", "", raw)

    try:
        slides = json.loads(raw)
        print(f"   ✅ Got {len(slides)} slides from LLM")
        return slides
    except json.JSONDecodeError:
        # Fallback: try to extract JSON array from anywhere in the text
        match = re.search(r"\[.*\]", raw, re.DOTALL)
        if match:
            slides = json.loads(match.group())
            print(f"   ✅ Extracted {len(slides)} slides (after cleanup)")
            return slides
        raise ValueError(f"LLM did not return valid JSON.\n\nRaw output:\n{raw[:500]}")


# ═══════════════════════════════════════════════════════════
# STEP 2 — Unsplash: Download one image per slide
# ═══════════════════════════════════════════════════════════

def download_image(keyword: str, save_path: Path, retries: int = 3) -> bool:
    """Fetch a landscape photo from Unsplash by keyword. Returns True on success."""
    if UNSPLASH_ACCESS_KEY == "YOUR_UNSPLASH_ACCESS_KEY":
        print("   ⚠️  Unsplash key not set — skipping images")
        return False

    for attempt in range(retries):
        try:
            resp = requests.get(
                "https://api.unsplash.com/photos/random",
                params={
                    "query": keyword,
                    "orientation": "landscape",
                    "client_id": UNSPLASH_ACCESS_KEY,
                },
                timeout=10,
            )
            if resp.status_code == 429:
                print(f"   ⏳ Rate limited — waiting 60s ...")
                time.sleep(60)
                continue
            resp.raise_for_status()
            data = resp.json()
            img_url = data["urls"]["regular"]  # ~1080px wide

            img_resp = requests.get(img_url, timeout=15)
            img_resp.raise_for_status()

            # Save as JPEG
            img = Image.open(BytesIO(img_resp.content)).convert("RGB")
            img.save(save_path, "JPEG", quality=85)
            return True

        except Exception as e:
            print(f"   ⚠️  Image fetch failed (attempt {attempt+1}): {e}")
            time.sleep(2)

    return False


def download_all_images(slides: list[dict]) -> dict[int, Path]:
    """Download images for all slides. Returns {slide_index: image_path}."""
    IMAGES_DIR.mkdir(exist_ok=True)
    image_map = {}

    print(f"\n🖼️  Downloading images ({len(slides)} slides) ...")
    for i, slide in enumerate(slides):
        keyword = slide.get("image_keyword", slide["title"])
        path = IMAGES_DIR / f"slide_{i:03d}.jpg"

        if path.exists():
            print(f"   [{i+1}/{len(slides)}] Cached: {keyword}")
            image_map[i] = path
            continue

        print(f"   [{i+1}/{len(slides)}] Fetching: '{keyword}' ...", end=" ", flush=True)
        success = download_image(keyword, path)
        if success:
            print("✅")
            image_map[i] = path
        else:
            print("❌ (will use color background)")

        time.sleep(0.5)  # be polite to the API

    return image_map


# ═══════════════════════════════════════════════════════════
# STEP 3 — PPTX: Assemble the presentation
# ═══════════════════════════════════════════════════════════

def hex_to_rgb(hex_str: str) -> RGBColor:
    h = hex_str.lstrip("#")
    return RGBColor(int(h[0:2], 16), int(h[2:4], 16), int(h[4:6], 16))


def set_slide_background(slide, color: RGBColor):
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_text_box(slide, text, left, top, width, height,
                 font_size=18, bold=False, color=None,
                 align=PP_ALIGN.LEFT, italic=False):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = Pt(font_size)
    run.font.bold = bold
    run.font.italic = italic
    if color:
        run.font.color.rgb = color
    return txBox


def add_image_half(slide, image_path: Path, side: str = "right"):
    """Add image to left or right half of slide."""
    half_w = SLIDE_W / 2
    left = half_w if side == "right" else Emu(0)
    try:
        slide.shapes.add_picture(
            str(image_path), left, Emu(0), half_w, SLIDE_H
        )
    except Exception as e:
        print(f"      ⚠️  Could not add image: {e}")


def add_full_bleed_image(slide, image_path: Path, opacity_overlay=True):
    """Full-width background image with dark overlay."""
    try:
        slide.shapes.add_picture(
            str(image_path), Emu(0), Emu(0), SLIDE_W, SLIDE_H
        )
        # Dark overlay rectangle
        if opacity_overlay:
            shape = slide.shapes.add_shape(
                1,  # MSO_SHAPE_TYPE.RECTANGLE
                Emu(0), Emu(0), SLIDE_W, SLIDE_H
            )
            shape.fill.solid()
            shape.fill.fore_color.rgb = RGBColor(0x00, 0x00, 0x00)
            shape.fill.fore_color.theme_color = None
            shape.line.fill.background()
            # Simulate transparency by using a very dark navy instead of black
            shape.fill.fore_color.rgb = RGBColor(0x06, 0x15, 0x2B)
    except Exception as e:
        print(f"      ⚠️  Could not add full-bleed image: {e}")


def build_title_slide(slide, data: dict, image_path=None):
    set_slide_background(slide, CLR_BG_DARK)
    if image_path:
        add_full_bleed_image(slide, image_path)

    # Accent bar on left
    bar = slide.shapes.add_shape(1, Inches(0), Inches(2.5), Inches(0.08), Inches(2.5))
    bar.fill.solid()
    bar.fill.fore_color.rgb = CLR_ACCENT
    bar.line.fill.background()

    add_text_box(
        slide, data["title"],
        left=Inches(0.6), top=Inches(2.4),
        width=Inches(9), height=Inches(1.5),
        font_size=44, bold=True, color=CLR_TEXT_LITE,
        align=PP_ALIGN.LEFT
    )
    subtitle = data["bullets"][0] if data.get("bullets") else ""
    add_text_box(
        slide, subtitle,
        left=Inches(0.6), top=Inches(4.0),
        width=Inches(9), height=Inches(0.7),
        font_size=20, color=CLR_ACCENT,
        align=PP_ALIGN.LEFT
    )


def build_section_slide(slide, data: dict, image_path=None):
    set_slide_background(slide, CLR_BG_DARK)
    if image_path:
        add_full_bleed_image(slide, image_path)

    add_text_box(
        slide, "— SECTION —",
        left=Inches(0.6), top=Inches(2.2),
        width=Inches(11), height=Inches(0.5),
        font_size=12, color=CLR_ACCENT, italic=True
    )
    add_text_box(
        slide, data["title"],
        left=Inches(0.6), top=Inches(2.8),
        width=Inches(11), height=Inches(1.5),
        font_size=38, bold=True, color=CLR_TEXT_LITE
    )
    summary = data["bullets"][0] if data.get("bullets") else ""
    add_text_box(
        slide, summary,
        left=Inches(0.6), top=Inches(4.4),
        width=Inches(11), height=Inches(0.8),
        font_size=16, color=RGBColor(0xCA, 0xDC, 0xFC)
    )


def build_content_slide(slide, data: dict, image_path=None):
    """Two-column layout: text left, image right."""
    set_slide_background(slide, CLR_BG_LIGHT)

    # Image on right half
    if image_path:
        add_image_half(slide, image_path, side="right")
        # Subtle gradient overlay on left edge of image (thin rect)
        fade = slide.shapes.add_shape(1, Inches(6.4), Emu(0), Inches(0.3), SLIDE_H)
        fade.fill.solid()
        fade.fill.fore_color.rgb = CLR_BG_LIGHT
        fade.line.fill.background()

    # Accent top bar
    bar = slide.shapes.add_shape(1, Emu(0), Emu(0), Inches(6.5), Inches(0.07))
    bar.fill.solid()
    bar.fill.fore_color.rgb = CLR_ACCENT
    bar.line.fill.background()

    # Title
    add_text_box(
        slide, data["title"],
        left=Inches(0.4), top=Inches(0.25),
        width=Inches(6.0), height=Inches(1.0),
        font_size=28, bold=True, color=CLR_TEXT_DARK
    )

    # Bullet points
    bullets = data.get("bullets", [])
    top = Inches(1.4)
    for bullet in bullets:
        # Bullet dot
        dot = slide.shapes.add_shape(
            1,
            Inches(0.4), top + Inches(0.12),
            Inches(0.12), Inches(0.12)
        )
        dot.fill.solid()
        dot.fill.fore_color.rgb = CLR_ACCENT
        dot.line.fill.background()

        add_text_box(
            slide, bullet,
            left=Inches(0.65), top=top,
            width=Inches(5.7), height=Inches(0.5),
            font_size=16, color=CLR_TEXT_DARK
        )
        top += Inches(0.72)


def build_conclusion_slide(slide, data: dict, image_path=None):
    set_slide_background(slide, CLR_BG_DARK)
    if image_path:
        add_full_bleed_image(slide, image_path)

    add_text_box(
        slide, data["title"],
        left=Inches(1), top=Inches(2.0),
        width=Inches(11), height=Inches(1.4),
        font_size=40, bold=True, color=CLR_TEXT_LITE,
        align=PP_ALIGN.CENTER
    )
    for i, bullet in enumerate(data.get("bullets", [])):
        add_text_box(
            slide, f"✦  {bullet}",
            left=Inches(2), top=Inches(3.6 + i * 0.65),
            width=Inches(9), height=Inches(0.55),
            font_size=16, color=CLR_ACCENT,
            align=PP_ALIGN.CENTER
        )


def assemble_presentation(slides: list[dict], image_map: dict) -> Presentation:
    print(f"\n📊 Assembling {len(slides)}-slide presentation ...")
    prs = Presentation()
    prs.slide_width  = SLIDE_W
    prs.slide_height = SLIDE_H

    blank_layout = prs.slide_layouts[6]  # completely blank

    builders = {
        "title":      build_title_slide,
        "section":    build_section_slide,
        "content":    build_content_slide,
        "conclusion": build_conclusion_slide,
    }

    for i, data in enumerate(slides):
        slide_type = data.get("slide_type", "content")
        builder = builders.get(slide_type, build_content_slide)
        image_path = image_map.get(i)

        slide = prs.slides.add_slide(blank_layout)
        builder(slide, data, image_path)

        print(f"   [{i+1}/{len(slides)}] {slide_type.upper()}: {data['title']}")

    return prs


# ═══════════════════════════════════════════════════════════
# MAIN
# ═══════════════════════════════════════════════════════════

def main():
    print("=" * 55)
    print("  🎯 AI PPT Generator — Powered by Ollama + Unsplash")
    print("=" * 55)

    topic      = input("\nEnter presentation topic: ").strip()
    num_slides = int(input("Number of slides (e.g. 75): ").strip())

    if not topic:
        print("❌ Topic cannot be empty.")
        return

    # Step 1: Generate slide content
    slides = generate_slide_plan(topic, num_slides)

    # Step 2: Download images
    image_map = download_all_images(slides)

    # Step 3: Build PPTX
    prs = assemble_presentation(slides, image_map)

    # Step 4: Save
    prs.save(OUTPUT_FILE)
    print(f"\n✅ Saved: {OUTPUT_FILE}  ({len(slides)} slides)")
    print("   Open it in PowerPoint or LibreOffice Impress.\n")


if __name__ == "__main__":
    main()
