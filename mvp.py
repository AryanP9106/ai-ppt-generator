import ollama
from pptx import Presentation

# 1. Connect to local Gemma 4
response = ollama.chat(model='gemma4', messages=[
    {'role': 'user', 'content': 'Give me a title and one bullet point about "Software Engineering Basics"'}
])

content = response['message']['content']

# 2. Create a basic PPT
prs = Presentation()
slide = prs.slides.add_slide(prs.slide_layouts[1])
slide.shapes.title.text = "AI Presentation"
slide.placeholders[1].text = content

# 3. Save
prs.save('test.pptx')
print("Successfully generated test.pptx")